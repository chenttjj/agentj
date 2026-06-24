#!/usr/bin/env python3
"""Build Mood Beat 專題報告.pptx with sage-terracotta theme."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

REPORT_DIR = Path("report")
ASSETS = REPORT_DIR / "assets"
OUT = REPORT_DIR / "專題報告.pptx"

# sage-terracotta palette
SAGE = RGBColor(0x87, 0x96, 0x73)
TERRACOTTA = RGBColor(0xC8, 0x6A, 0x4F)
CREAM = RGBColor(0xF5, 0xEF, 0xE4)
DARK = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def add_text(slide, left, top, width, height, text, size=20, color=DARK, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Microsoft JhengHei"
    return tb


def add_accent(slide, color, height):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_footer(slide, page, total):
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.7), Inches(0.3),
             f"{page} / {total}", size=10, color=DARK)


def title_slide(page, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_accent(s, TERRACOTTA, Inches(0.6))
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
             "Mood Beat", size=60, color=TERRACOTTA, bold=True)
    add_text(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.8),
             "心情 × 呼吸 × 音樂 — 你的專屬節奏", size=28, color=SAGE, bold=True)
    add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
             "馬公高中 Agent Studio 專題報告", size=18, color=DARK)
    add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
             "完成日期：2026-06-24", size=14, color=DARK)
    add_accent(s, SAGE, Inches(0.2))
    add_footer(s, page, total)


def content_slide(title, bullets, page, total, accent=SAGE):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_accent(s, accent, Inches(0.4))
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             title, size=32, color=TERRACOTTA, bold=True)
    y = Inches(1.9)
    for b in bullets:
        add_text(s, Inches(0.8), y, Inches(11.5), Inches(0.6),
                 "• " + b, size=20, color=DARK)
        y = Emu(int(y) + int(Inches(0.65)))
    add_accent(s, accent, Inches(0.15))
    add_footer(s, page, total)


def image_slide(title, img, caption, page, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_accent(s, SAGE, Inches(0.4))
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.6),
             title, size=28, color=TERRACOTTA, bold=True)
    s.shapes.add_picture(str(img), Inches(1.5), Inches(1.6), height=Inches(5.0))
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             caption, size=12, color=DARK)
    add_accent(s, TERRACOTTA, Inches(0.15))
    add_footer(s, page, total)


TOTAL = 8

title_slide(1, TOTAL)
content_slide("01  專題介紹", [
    "專題名稱：Mood Beat",
    "一句話：讓使用者描述今天狀態，透過呼吸穩定節奏，最後產生專屬 beat",
    "主要使用者：喜歡音樂且喜歡心靈呼吸放鬆的人",
    "解決痛點：壓力大、心情不愉快",
], 2, TOTAL, TERRACOTTA)
image_slide("02  學校 Server 環境", ASSETS / "server-topology.png",
            "圖：全班相同 server 拓撲圖（不揭露 IP／URL）", 3, TOTAL)
image_slide("03  Mood Beat 個人架構", ASSETS / "project-architecture.png",
            "圖：左欄 Mood Beat 頁 → 右欄 Agent → JSON 節拍參數", 4, TOTAL)
content_slide("04  左欄 ↔ 右欄 互動", [
    "左欄 Mood Beat：心情值滑桿、壓力值滑桿、音樂風格下拉",
    "按鈕：「開始呼吸」、「生成 Beat」",
    "傳給 Agent：心情值、壓力值、音樂偏好",
    "Agent 回傳：文字建議 + JSON 節拍參數",
    "完整例子：心情 7 / 壓力 3 / lo-fi → 「你的專屬 lo-fi beat 已生成」",
], 5, TOTAL, SAGE)
content_slide("05  成果", [
    "用滑桿輸入心情值與壓力值",
    "選擇音樂風格（lo-fi、trap、輕音樂）",
    "按鈕啟動呼吸引導",
    "根據情緒生成音樂風格建議",
    "Agent 回傳文字建議 + JSON 節拍參數",
], 6, TOTAL, TERRACOTTA)
content_slide("06  創新與技術", [
    "亮點：把「情緒紀錄、呼吸穩定、音樂推薦」三步驟結合",
    "亮點：依心情／壓力／偏好產出個人化 beat 建議",
    "亮點：先穩定心情、再創作音樂的紓壓體驗",
    "技術：Agent Studio + Streamlit（左欄頁面）",
    "技術：Agent.chat（右欄對話與推理）",
    "技術：JSON 傳遞情緒資料、LLM 生成音樂風格建議",
], 7, TOTAL, SAGE)
image_slide("07  Demo 截圖 — Mood Beat 頁面", ASSETS / "demo-02.png",
            "圖：Mood Beat 自訂頁實際畫面", 8, TOTAL)

prs.save(OUT)
print(f"OK: {OUT}")
